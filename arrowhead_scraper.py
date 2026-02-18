#!/usr/bin/env python3
"""
Arrowhead Pharmaceuticals – Events & Presentations Scraper + JSON Extractor
Version 4 — curl-cffi Akamai TLS bypass
============================================================================

WHY curl-cffi?
──────────────
The site uses Akamai Bot Manager which performs TLS fingerprinting (JA3/JA4).
Standard Python libraries (requests, urllib3, httpx) have a distinct TLS
handshake that Akamai recognises and resets with HTTP/2 INTERNAL_ERROR.

curl-cffi embeds libcurl compiled with BoringSSL, impersonating Chrome's
exact TLS fingerprint. Akamai cannot distinguish it from a real browser.

SETUP:
──────
    pip install curl-cffi beautifulsoup4 lxml pdfplumber python-pptx

That's it — no cookies file, no Selenium, no proxy needed.

HOW IT WORKS:
─────────────
  - Fetches the base events page (current year loaded by default).
  - For prior years, submits the year-filter GET form with the correct params.
  - Extracts all /static-files/ PDF/PPTX links via BeautifulSoup.
  - Downloads each new file, extracts content to JSON, deletes the original.
  - Tracks processed URLs in arrowhead_seen.json to avoid re-downloading.
"""

import re, json, time, logging
from datetime import datetime, timezone
from pathlib import Path
from urllib.parse import urljoin

from curl_cffi import requests          # drop-in replacement, Chrome TLS fingerprint
from bs4 import BeautifulSoup
import pdfplumber
from pptx import Presentation

# ── Config ─────────────────────────────────────────────────────────────────

BASE_URL        = "https://ir.arrowheadpharma.com/events-and-presentations"
DOMAIN          = "https://ir.arrowheadpharma.com"
DOWNLOAD_DIR    = Path("./arrowhead_downloads")
STATE_FILE      = Path("./arrowhead_seen.json")
LOOKBACK_MONTHS = 36
KEEP_ORIGINALS  = False
PAGE_DELAY      = 3    # seconds between page requests
FILE_DELAY      = 1    # seconds between file downloads

# Found in the <select name="HASH_year[value]"> on the events page.
# Update this constant if the site ever rebuilds and the hash changes.
WIDGET_HASH = "a55e2d9e49ef1b96eb6f3e69698dc5fce32d89ec43e78829900575cec31a2523"

HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
        "AppleWebKit/537.36 (KHTML, like Gecko) "
        "Chrome/132.0.0.0 Safari/537.36"
    ),
    "Accept":                    "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
    "Accept-Language":           "en-US,en;q=0.9",
    "Accept-Encoding":           "gzip, deflate, br",
    "Connection":                "keep-alive",
    "Upgrade-Insecure-Requests": "1",
    "Sec-Fetch-Dest":            "document",
    "Sec-Fetch-Mode":            "navigate",
    "Sec-Fetch-Site":            "none",
}

FILE_EXTS = {".pdf", ".ppt", ".pptx"}

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)
log = logging.getLogger(__name__)

# ══════════════════════════════════════════════════════════════════════════
# SESSION  (curl-cffi — Chrome TLS fingerprint bypasses Akamai JA3/JA4)
# ══════════════════════════════════════════════════════════════════════════

def build_session() -> requests.Session:
    """
    curl-cffi's Session with impersonate="chrome124" sends an authentic
    Chrome TLS ClientHello. Akamai Bot Manager cannot distinguish it from
    a real browser, so the HTTP/2 INTERNAL_ERROR reset never occurs.
    """
    s = requests.Session(impersonate="chrome124")
    s.headers.update(HEADERS)
    return s


# ══════════════════════════════════════════════════════════════════════════
# FETCHING
# ══════════════════════════════════════════════════════════════════════════

def fetch(url: str, session: requests.Session, params: dict | None = None,
          referer: str | None = None) -> str | None:
    hdrs = dict(HEADERS)
    if referer:
        hdrs["Referer"]        = referer
        hdrs["Sec-Fetch-Site"] = "same-origin"
    try:
        r = session.get(url, params=params, headers=hdrs, timeout=45)
        r.raise_for_status()
        log.info("  HTTP %d  %d bytes  %s", r.status_code, len(r.content),
                 (url + ("?" + "&".join(f"{k}={v}" for k, v in (params or {}).items()))
                  if params else url)[:120])
        return r.text
    except Exception as exc:
        log.error("  Request failed: %s — %s", url, exc)
        return None


def target_years() -> list[int]:
    now = datetime.now()
    earliest = now.year - (LOOKBACK_MONTHS // 12)
    if (now.month - (LOOKBACK_MONTHS % 12)) <= 0:
        earliest -= 1
    return list(range(earliest, now.year + 1))


def load_seen() -> set:
    if STATE_FILE.exists():
        try:
            return set(json.loads(STATE_FILE.read_text()))
        except Exception:
            pass
    return set()


def save_seen(seen: set) -> None:
    STATE_FILE.write_text(json.dumps(sorted(seen), indent=2))


def extract_links(html: str) -> list[dict]:
    soup  = BeautifulSoup(html, "lxml")
    found, seen_urls = [], set()
    for tag in soup.find_all("a", href=True):
        href     = tag["href"].strip()
        full_url = urljoin(DOMAIN, href) if href.startswith("/") else href
        if "/static-files/" not in full_url or full_url in seen_urls:
            continue
        filename = (tag.get("title") or "").strip() or tag.get_text(strip=True) or Path(full_url).name
        if Path(filename).suffix.lower() in FILE_EXTS:
            seen_urls.add(full_url)
            found.append({"url": full_url, "filename": filename})
    return found


def safe_name(name: str) -> str:
    return re.sub(r'[\\/*?:"<>|]', "_", name).strip()


def download_file(url: str, filename: str, session: requests.Session) -> Path | None:
    DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)
    filename  = safe_name(filename)
    dest      = DOWNLOAD_DIR / filename
    try:
        # curl_cffi doesn't support streaming context manager — fetch whole response
        r = session.get(url, headers={**HEADERS, "Referer": BASE_URL}, timeout=120)
        r.raise_for_status()
        cd = r.headers.get("Content-Disposition", "")
        m  = re.search(r'filename=["\']?([^"\';\n]+)', cd)
        if m:
            dest = DOWNLOAD_DIR / safe_name(m.group(1).strip())
        stem, suffix, n = dest.stem, dest.suffix, 1
        while dest.exists():
            dest = DOWNLOAD_DIR / f"{stem}_{n}{suffix}"; n += 1
        dest.write_bytes(r.content)
        log.info("    ✓ Downloaded: %s (%d KB)", dest.name, len(r.content) // 1024)
        return dest
    except Exception as exc:
        log.error("    ✗ Download failed: %s", exc)
        return None


# ══════════════════════════════════════════════════════════════════════════
# PPTX EXTRACTION
# ══════════════════════════════════════════════════════════════════════════

_NS = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"}

def _pptx_chart(chart) -> dict:
    r = {"type": "chart", "chart_type": "unknown", "title": None,
         "categories": [], "series": []}
    try:
        r["chart_type"] = chart.chart_type.name
    except Exception:
        pass
    try:
        if chart.has_title and chart.chart_title.has_text_frame:
            r["title"] = chart.chart_title.text_frame.text.strip()
    except Exception:
        pass
    try:
        r["categories"] = [n.text for n in chart._element.findall(".//c:cat//c:v", _NS) if n.text]
    except Exception:
        pass
    try:
        for s in chart.series:
            entry = {"name": getattr(s, "name", None), "values": []}
            try:
                entry["values"] = [round(v, 6) if isinstance(v, float) else v for v in s.values]
            except Exception:
                entry["values"] = [n.text for n in s._element.findall(".//c:v", _NS) if n.text]
            r["series"].append(entry)
    except Exception:
        pass
    return r

def extract_pptx(path: Path, url: str) -> dict:
    prs = Presentation(str(path))
    doc = {"metadata": {"source_filename": path.name, "source_url": url,
                        "scraped_at": datetime.now(timezone.utc).isoformat(),
                        "file_size_bytes": path.stat().st_size,
                        "format": "pptx", "slide_count": len(prs.slides)},
           "slides": []}
    for i, slide in enumerate(prs.slides, 1):
        entry = {"slide_number": i, "title": None, "content": []}
        try:
            if slide.shapes.title:
                entry["title"] = slide.shapes.title.text.strip() or None
        except Exception:
            pass
        for shape in slide.shapes:
            if shape.has_chart:
                entry["content"].append(_pptx_chart(shape.chart)); continue
            try:
                if shape.has_table:
                    entry["content"].append({"type": "table",
                        "rows": [[c.text.strip() for c in r.cells] for r in shape.table.rows]}); continue
            except Exception:
                pass
            try:
                if shape.has_text_frame:
                    lines = [" ".join(r.text for r in p.runs).strip()
                             for p in shape.text_frame.paragraphs]
                    txt = "\n".join(l for l in lines if l)
                    if txt and txt != (entry["title"] or ""):
                        entry["content"].append({"type": "text", "text": txt})
            except Exception:
                pass
        doc["slides"].append(entry)
    return doc


# ══════════════════════════════════════════════════════════════════════════
# PDF EXTRACTION
# ══════════════════════════════════════════════════════════════════════════

def _is_chart_table(tbl):
    if not tbl or len(tbl) < 2:
        return False
    cells = [c for row in tbl[1:] for c in row[1:]]
    if not cells:
        return False
    return sum(1 for c in cells if re.match(r'^-?[\d,\.%]+$', c.strip())) / len(cells) >= 0.6

def _tbl_to_chart(tbl):
    cats = tbl[0][1:]
    series = []
    for row in tbl[1:]:
        if not row: continue
        vals = []
        for c in row[1:]:
            try: vals.append(float(re.sub(r'[,$%]', '', c.strip())))
            except: vals.append(c)
        series.append({"name": row[0], "values": vals})
    return {"type": "chart", "chart_type": "table_derived",
            "title": None, "categories": cats, "series": series}

def extract_pdf(path: Path, url: str) -> dict:
    doc = {"metadata": {"source_filename": path.name, "source_url": url,
                        "scraped_at": datetime.now(timezone.utc).isoformat(),
                        "file_size_bytes": path.stat().st_size,
                        "format": "pdf", "page_count": 0},
           "pages": []}
    with pdfplumber.open(str(path)) as pdf:
        doc["metadata"]["page_count"] = len(pdf.pages)
        for i, page in enumerate(pdf.pages, 1):
            entry = {"page_number": i, "content": []}
            for raw in (page.extract_tables() or []):
                tbl = [[c or "" for c in r] for r in raw if r]
                if tbl:
                    entry["content"].append(_tbl_to_chart(tbl) if _is_chart_table(tbl)
                                            else {"type": "table", "rows": tbl})
            txt = re.sub(r'\n{3,}', '\n\n', page.extract_text(layout=True) or "").strip()
            if txt:
                entry["content"].append({"type": "text", "text": txt})
            doc["pages"].append(entry)
    return doc


# ══════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════

def main():
    log.info("=== Arrowhead scraper v4 (curl-cffi, lookback: %d months) ===", LOOKBACK_MONTHS)

    session = build_session()
    seen    = load_seen()
    years   = target_years()
    log.info("Targeting years: %s", years)

    # ── Step 1: Base page (no params — current year loaded by default) ────
    log.info("── Fetching base page...")
    base_html = fetch(BASE_URL, session)

    if not base_html:
        log.error("Base page fetch failed. Check your internet connection.")
        return

    all_links = extract_links(base_html)
    log.info("  Base page: %d file link(s) found", len(all_links))

    # ── Step 2: Per-year filtered pages ──────────────────────────────────
    current_year = datetime.now().year
    for year in [y for y in years if y != current_year]:
        log.info("── Year: %d", year)
        time.sleep(PAGE_DELAY)
        params = {
            f"{WIDGET_HASH}_year[value]": str(year),
            f"{WIDGET_HASH}_widget_id":   WIDGET_HASH,
            "op":      "Filter",
            "form_id": "widget_form_base",
        }
        html = fetch(BASE_URL, session, params=params, referer=BASE_URL)
        if not html:
            log.warning("  Year %d skipped (blocked or timeout).", year)
            continue
        links = extract_links(html)
        log.info("  %d file link(s) found.", len(links))
        all_links.extend(links)

    unique = list({item["url"]: item for item in all_links}.values())
    log.info("Total unique links: %d", len(unique))

    # ── Step 3: Download + Extract ────────────────────────────────────────
    new_count = 0
    for item in unique:
        url = item["url"]
        if url in seen:
            continue
        log.info("  ↓ %s", item["filename"])
        time.sleep(FILE_DELAY)
        path = download_file(url, item["filename"], session)
        if not path:
            continue
        ext  = path.suffix.lower()
        try:
            data = extract_pdf(path, url) if ext == ".pdf" else extract_pptx(path, url)
            jp   = path.with_suffix(".json")
            jp.write_text(json.dumps(data, indent=2, ensure_ascii=False))
            log.info("    ✓ JSON: %s (%d KB)", jp.name, jp.stat().st_size // 1024)
        except Exception as exc:
            log.error("    ✗ Extraction error: %s", exc)
            seen.add(url); save_seen(seen); continue
        if not KEEP_ORIGINALS:
            try:
                path.unlink()
                log.info("    🗑  Deleted: %s", path.name)
            except OSError:
                pass
        seen.add(url)
        new_count += 1

    save_seen(seen)
    log.info("Done. %d new file(s) processed.", new_count)
    log.info("=== Arrowhead scraper finished ===\n")


if __name__ == "__main__":
    main()